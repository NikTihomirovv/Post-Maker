import base64
import logging
import os
import tempfile
from dataclasses import dataclass, field, is_dataclass
from enum import Enum
from pprint import pprint
from typing import Any, Dict, List, Optional

import cairosvg
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class PresentationManager:
    """Менеджер для работы с презентациями."""

    def __init__(self, settings):
        self._logger = logging.getLogger(__name__)
        self._settings = settings
        self._presentation_factory = _PresentationFactory()
        

    def create(self, _data: PresentationData):

        kwargs = {'settings': self._settings}

        self._presentation = self._presentation_factory._create(**kwargs)
        self._presentation._create(_data)


    




class _Presentation:
    """Класс для описания презентации."""

    def __init__(self, **kwargs):
        self._logger = logging.getLogger(__name__)
        self._settings = kwargs.get('settings')

        self._slide_layout_factory = _SlideLayoutFactory()

        self._name = None
        self._width = None
        self._height = None
        self._slides = None


    def _create(self, _data: PresentationData)  -> None:
        """Создание объекта презентации"""
    
        try:
            # Получаем объекты PresentationData
            if _data and is_dataclass(_data):
                self._logger.debug(f'✅ Получен объект презентации для обработки')

                self._name = _data.name if isinstance(_data.name, str) and _data.name != '' else 'Без имени'
                self._prs_width = _data.width if isinstance(_data.width, int) and _data.width > 0 else 1920
                self._prs_height = _data.height if isinstance(_data.height, int) and _data.height > 0 else 1080
                self._slides = _data.slides if isinstance(_data.slides, list) and len(_data.slides) != 0 else []

                width_emu = int(self._prs_width * 12700) if self._prs_width > 0 else 1920 * 12700
                height_emu = int(self._prs_height * 12700) if self._prs_height > 0 else 1080 * 12700

                _prs = Presentation()
                _prs.slide_width = width_emu
                _prs.slide_height = height_emu

                self._logger.debug(f'   ✅ Размер слайда установлен:')
                self._logger.debug(f'   ✅ EMU: {width_emu} x {height_emu}')
                self._logger.debug(f'   ✅ Пиксели: {self._prs_width} x {self._prs_height} px')
            
                self._get_prs_slides(self._slides, _prs, self._prs_width, self._prs_height)
                _prs.save('test.pptx')

            else:
                self._logger.error(f'❌ Ошибка в получении презентации для обработке')
     
        except Exception as e:
            self._logger.error(f'❌ Ошибка при создании презентации: {e}')


    def _get_prs_slides(self, _slides: List[_Slide], _prs: Presentation, _prs_width: int, _prs_height: int) -> None:
        
        if _slides:
            self._logger.debug(f'✅ Получено {len(_slides)} слайдов для обработки')

            for idx, _slide in enumerate(_slides, 1):
        
                if not is_dataclass(_slide):
                    self._logger.error(f'❌ Ошибка в обработке {idx} слайда')
                    continue
        
                # Получаем объекты Slide
                _slide_layout_type = _slide.layout if _slide.layout else None
                _slide_objects = _slide.objects if _slide.objects else None
        
                if not isinstance(_slide_layout_type, Enum):
                    self._logger.error(f'❌ Ошибка в обработке типа лэйаута {_slide_layout_type}')
                    continue
        
                if not isinstance(_slide_objects, list) or len(_slide_objects) < 1:
                    self._logger.error(f'❌ Ошибка в обработке объектов слайда {_slide_objects}')
                    continue

                # Создаем слайды
                kwargs = {
                    'settings': self._settings,
                    'prs': _prs,
                    'slide_layout_type': _slide_layout_type,
                    'slide_objects': _slide_objects,
                    'prs_width': _prs_width,
                    'prs_height': _prs_height
                    }
                _slide_layout = self._slide_layout_factory._create(**kwargs)
                _slide_layout._create()
            
        else:
            self._logger.warning(f'❌ Отсутствуют слайды для обработке')
        


class _CustomSlideLayout:
    """Класс для описания слайда с кастомным лэйаутом."""

    def __init__(self, **kwargs):
        self._logger = logging.getLogger(__name__)

        self._object_factory = _ObjectFactory()

        self._settings = kwargs.get('settings')
        self._prs = kwargs.get('prs')
        self._prs_width = kwargs.get('prs_width')
        self._prs_height = kwargs.get('prs_height')
        self._slide_layout_type = kwargs.get('slide_layout_type')
        self._slide_objects = kwargs.get('slide_objects')


    def _create(self) -> None:

        try: 
            if self._slide_layout_type:
                self._logger.debug(f'✅ Создаем слайд с типом лэйаута {self._slide_layout_type.name}')
                _slide_layout = self._prs.slide_layouts[self._slide_layout_type.value]
                _slide = self._prs.slides.add_slide(_slide_layout)

                if _slide: 
                    for _slide_object in self._slide_objects:

                        if is_dataclass(_slide_object):
                            self._logger.debug(f'   ✅ Получен объект слайда')

                            _object_type = _slide_object.type
                            if not isinstance(_object_type, Enum):
                                self._logger.error(f'❌ Ошибка в получении типа объекта')
                                continue

                            _object_text_style = _slide_object.text_style
                            if _object_text_style != None:
                                if not is_dataclass(_object_text_style):
                                    self._logger.error(f'❌ Ошибка в получении стилей объекта')
                                    continue
                    
                            _object_content = _slide_object.content
                            if not _object_content:
                                self._logger.error(f'❌ Ошибка в получении контента объекта')
                                continue

                            _object_position = _slide_object.position
                            if not isinstance(_object_position, tuple) or len(_object_position) != 2:
                                self._logger.error(f'❌ Ошибка в получении позиции объекта')
                                continue

                            if not all(isinstance(coord, (int, float)) for coord in _object_position):
                                self._logger.error(f'❌ Координаты должны быть числами, получено {_object_position}')
                                continue

                            _object_size = _slide_object.size
                            if not isinstance(_object_size, tuple) or len(_object_size) != 2:
                                self._logger.error(f'❌ Ошибка в получении размера объекта')
                                continue

                            if not all(isinstance(size, (int, float)) for size in _object_size):
                                self._logger.error(f'❌ Размеры должны быть числами, получено {_object_position}')
                                continue

                            _vertically_centered = _slide_object.vertically_centered
                            if not isinstance(_vertically_centered, bool):
                                self._logger.error(f'❌ Ошибка в получении параметров центрирования')
                                continue

                            _horizontally_centred = _slide_object.horizontally_centred
                            if not isinstance(_horizontally_centred, bool):
                                self._logger.error(f'❌ Ошибка в получении параметров центрирования')
                                continue
                            
                            _object_z_index = _slide_object.z_index
                            if not isinstance(_object_z_index, int):
                                self._logger.error(f'❌ Ошибка в получении z-index объекта')
                                continue

                            _x, _y = _object_position
                            _w, _h = _object_size

                            # Создаем объекты на слайде
                            kwargs = {
                                'settings': self._settings,
                                'prs_width': self._prs_width,
                                'prs_height': self._prs_height,
                                'slide': _slide,
                                'type': _object_type,
                                'text_style': _object_text_style,
                                'content': _object_content,
                                'x': _x,
                                'y': _y,
                                'w': _w,
                                'h': _h,
                                'vertically_centered': _vertically_centered,
                                'horizontally_centred': _horizontally_centred,
                                'z_index': _object_z_index,
                            }
                            _object = self._object_factory._create(**kwargs)
                            _object._create()

                        else:
                            self._logger.error(f'❌ Ошибка в получении объекта слайда')

            else:
                self._logger.error(f'❌ Ошибка в применении слайд лэйаута')
                return False

        except Exception as e:
            self._logger.error(f'❌ Ошибка при создании слайда: {e}')



# Объекты на слайде ======================================================================================


class _BaseObj:
    """Базовый класс для создания объектов на слайде"""

    def __init__(self, **kwargs):
        self._logger = logging.getLogger(__name__)

        self._settings = kwargs.get('settings')
        self._prs_width = kwargs.get('prs_width')
        self._prs_height = kwargs.get('prs_height')
        self._slide = kwargs.get('slide')
        self._type = kwargs.get('type')
        self._text_style = kwargs.get('text_style', None)
        self._content = kwargs.get('content')
        
        self._x = self._px_to_emu(kwargs.get('x', 0))
        self._y = self._px_to_emu(kwargs.get('y', 0))
        self._w = self._px_to_emu(kwargs.get('w', 100))
        self._h = self._px_to_emu(kwargs.get('h', 100))
        self._vertically_centered = kwargs.get('vertically_centered', False)
        self._horizontally_centred = kwargs.get('horizontally_centred', False)
        self._z_index = kwargs.get('z_index')
    
    @staticmethod
    def _px_to_emu(px: int) -> int:
        """Конвертирует пиксели в EMU (1 px = 12700 EMU)"""
        return int(px * 12700)
    
    @staticmethod
    def _emu_to_px(emu: int) -> int:
        """Конвертирует EMU в пиксели"""
        return int(emu / 12700)


    def _apply_text_style(self, _text_frame) -> None:
        """Применяет стиль к тексту"""

        try:
            if not _text_frame or not self._text_style:
                return
            
            # Получаем первый параграф
            if len(_text_frame.paragraphs) == 0:
                _text_frame.paragraphs.add()
            
            paragraph = _text_frame.paragraphs[0]
            
            # Применяем все стили
            if self._text_style.font_size:
                paragraph.font.size = int(self._text_style.font_size * 12700)  # Убедимся что int
            
            if self._text_style.font_name:
                paragraph.font.name = self._text_style.font_name
                
            if self._text_style.bold is not None:
                paragraph.font.bold = self._text_style.bold
                
            if self._text_style.italic is not None:
                paragraph.font.italic = self._text_style.italic
            
            # Применяем цвет
            if hasattr(self._text_style, 'color') and self._text_style.color:
                color_str = self._text_style.color
                if color_str.startswith('#') and len(color_str) == 7:
                    r = int(color_str[1:3], 16)
                    g = int(color_str[3:5], 16)
                    b = int(color_str[5:7], 16)
                    paragraph.font.color.rgb = RGBColor(r, g, b)
                    self._logger.debug(f'      ✅ Применен цвет: {color_str}')

            
            # Выравнивание
            if hasattr(self._text_style, 'alignment') and self._text_style.alignment:
                alignment_map = {
                    'left': PP_ALIGN.LEFT,
                    'center': PP_ALIGN.CENTER,
                    'right': PP_ALIGN.RIGHT,
                }
                if self._text_style.alignment in alignment_map:
                    paragraph.alignment = alignment_map[self._text_style.alignment]
                        
        except Exception as e:
            self._logger.error(f'   ❌ Ошибка применения стиля: {e}')


    def _check_alignment(self) -> None:
        _prs_width_emu = self._px_to_emu(self._prs_width)
        _prs_height_emu = self._px_to_emu(self._prs_height)

        if self._vertically_centered:
            self._y = (_prs_height_emu - self._h) // 2
            self._logger.debug(f'      ✅ Центрируем по вертикали: y={self._y} EMU')
        
        if self._horizontally_centred:
            self._x = (_prs_width_emu - self._w) // 2
            self._logger.debug(f'      ✅ Центрируем по горизонтали: x={self._x} EMU')


class _Title(_BaseObj):
    """Класс для описания заголовка."""

    def _create(self) -> None:
        try:
            self._logger.debug('      ✅ Создаем Title')

            self._check_alignment()

            _textbox = self._slide.shapes.add_textbox(
                self._x, self._y, self._w, self._h
            )
            
            _text_frame = _textbox.text_frame
            _text_frame.text = str(self._content)
            
            # Применяем стили
            if self._text_style:
                self._apply_text_style(_text_frame)
                
            self._logger.debug(f'      ✅ Title создан как текстовый блок')

        except Exception as e:
            self._logger.error(f'❌ Ошибка при создании Title: {e}')


class _SubTitle(_BaseObj):
    """Класс для описания подзаголовка."""


class _Text(_BaseObj):
    """Класс для описания текста."""


class _Image(_BaseObj):
    """Класс для описания картинки."""

    def _create(self) -> None:

        try:
            self._logger.debug('      ✅ Создаем Image')
            
            if self._content.startswith('data:image/png;base64,'):
                self._logger.debug('      ✅ Получено изображение в base64')

                _base64_data = self._content.split(',', 1)[1]

                # Декодируем base64 в байты
                _image_bytes = base64.b64decode(_base64_data)

                # Создаем временный PNG файл
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as _tmp_file:
                    _tmp_file.write(_image_bytes)
                    _tmp_path = _tmp_file.name

                # Получаем реальные размеры изображения в пикселях
                with PILImage.open(_tmp_path) as img:
                    img_width_px, img_height_px = img.size
                    self._logger.debug(f'      ✅ Размер изображения: {img_width_px}x{img_height_px} px')

                # Добавляем изображение
                self._slide.shapes.add_picture(
                    _tmp_path,
                    self._x, self._y,
                    self._w, self._h
                )
                
                # Удаляем временный файл
                os.unlink(_tmp_path)
                
                self._logger.info(f'      ✅ Image добавлено')

            else:
                self._logger.debug('      ❌ Неверный формат изображения')
            
        except Exception as e:
            self._logger.error(f'❌ Ошибка при создании Image: {e}')
            try:
                os.unlink(_tmp_path)
            except:
                pass


class _Table(_BaseObj):
    """Класс для описания таблицы."""


class _Chart(_BaseObj):
    """Класс для описания чарта."""


class _Background(_BaseObj):
    """Класс для описания бэкграунда."""



# Фабрики ================================================================================================



class _PresentationFactory:
    """Фабрика для создания презентаций."""

    @staticmethod
    def _create(**kwargs):
        return _Presentation(**kwargs)



class _SlideLayoutFactory():
    """Фабрика для создания слайда с одним из дефолтных лэйаутов"""

    @staticmethod
    def _create(**kwargs):

        _slide_layout_type = kwargs.get('slide_layout_type')

        match _slide_layout_type.name:

            case 'BLANK':
                return _CustomSlideLayout(**kwargs)

            case _:
                pass
                


class _ObjectFactory:
    """Фабрика для создания объектов слайда"""

    @staticmethod
    def _create(**kwargs):
        _object_type = kwargs.get('type')

        match _object_type.value:

            case 'title':
                return _Title(**kwargs)

            case 'subtitle':
                return _SubTitle(**kwargs)

            case 'text':
                return _Text(**kwargs)

            case 'image':
                return _Image(**kwargs)

            case 'table':
                return _Table(**kwargs)

            case 'chart':
                return _Chart(**kwargs)

            case 'background':
                return _Background(**kwargs)



# ENUMS ==================================================================================================
class _SlideLayout(Enum):
    """Типы макетов слайдов PowerPoint"""
    
    TITLE = 0                      # Title (presentation title slide)
    TITLE_AND_CONTENT = 1          # Title and Content
    SECTION_HEADER = 2             # Section Header (sometimes called Segue)
    TWO_CONTENT = 3                # Two Content (side by side bullet textboxes)
    COMPARISON = 4                 # Comparison (same but additional title for each side by side content box)
    TITLE_ONLY = 5                 # Title Only
    BLANK = 6                      # Blank
    CONTENT_WITH_CAPTION = 7       # Content with Caption
    PICTURE_WITH_CAPTION = 8       # Picture with Caption


class _ObjectType(Enum):
    """Типы объектов на слайде"""
    TITLE = 'title'
    SUBTITLE = 'subtitle'
    TEXT = 'text'
    IMAGE = 'image'
    TABLE = 'table'
    CHART = 'chart'
    BACKGROUND = 'background'



# Датаклассы =============================================================================================
@dataclass
class _ObjectTextStyle:
    """Стиль объекта"""

    font_size: Optional[int] = None
    font_name: Optional[str] = None
    bold: bool = False
    italic: bool = False
    color: Optional[str] = None
    alignment: str = 'left'


@dataclass
class _SlideObject:
    """Объект на слайде"""

    type: _ObjectType 
    content: Any          # str для текста, bytes/path для изображений
    text_style: Optional[_ObjectTextStyle] = None

    size: tuple[int, int] = (100, 100) 
    position: tuple[int, int] = (0, 0)
    vertically_centered: bool = False
    horizontally_centred: bool = False
    z_index: Optional[int] = None


@dataclass
class _Slide:
    """Слайд презентации"""

    layout: _SlideLayout
    objects: List[_SlideObject] = field(default_factory=list)


@dataclass
class PresentationData:
    """Данные для создания презентации"""

    name: str
    slides: List[_Slide] = field(default_factory=list)
    width: int = 1920
    height: int = 1080



# Паттерны ===============================================================================================



class PresentationPatterns:
    """Класс для хранения паттернов."""
    
    def __init__(self, settings):
        self._logger = logging.getLogger(__name__)
        self._settings = settings


    def _load_svg(self, _file_path, width=1920, height=1080) -> str:
        """Загружает SVG и конвертирует в PNG с указанными размерами через CairoSVG, возвращает base64."""
        try:
            png_data = cairosvg.svg2png(
                url=_file_path,
                output_width=width,
                output_height=height
            )
            
            if not png_data:
                self._logger.error('❌ Ошибка: PNG пустой')
                return None
            
            base64_str = base64.b64encode(png_data).decode('utf-8')
            self._logger.debug('✅ Изображение получено')
            return f"data:image/png;base64,{base64_str}"
            
        except Exception as e:
            self._logger.error(f'❌ Ошибка загрузки SVG: {e}')
            return 'Не удалось получить изображение'


    def get_pattern(self, _pattern_name: str, _data: Dict[Any]) -> PresentationData:

        _data = {
            'presentation_name': '',

            # Первый слайд
            'slide_1_text_1': 'Влияние загрязнения воздуха на активность заболевания у пациентов с ревматоидным артритом',
            'slide_1_text_2': 'Sung Cheol Jung, Seojin Yang, Min Hyuk Lim, Saram Lee, Sung Ik Cho, Jin Kyun Park, Jun Won Park, Eun Bong Lee (Seoul National University College of Medicine, Республика Корея)',
            'slide_1_text_3': 'Annals of the Rheumatic Diseases (ARD), The EULAR Journal, Elsevier, 2026',
            'slide_1_text_4': 'Ревматоидный артрит (RA), загрязнение воздуха, PM2.5, воспаление, аутоиммунные заболевания',

            # Второй слайд
            'slide_2_text_1': 'Ревматоидный артрит — хроническое аутоиммунное заболевание, поражающее 0.5-1% взрослого населения мира. Курение уже признано основным фактором риска. Предыдущие исследования показали связь загрязнения воздуха с повышенной вероятностью развития РА.',
            'slide_2_text_2': 'Неизвестно, влияет ли загрязнение воздуха на активность заболевания и провоцирует ли обострения у людей, уже имеющих РА. Также не изучены биологические механизмы этой связи.',
            'slide_2_text_3': 'Определить, влияет ли загрязнение воздуха на активность заболевания и риск обострений у пациентов с ревматоидным артритом, а также исследовать биологические процессы, объясняющие эту связь.',
            'slide_2_text_4': 'Воздействие загрязненного воздуха, особенно мелкодисперсных частиц PM2.5, связано с повышенной активностью РА и увеличением частоты обострений через механизмы окислительного стресса и воспаления.',

            'slide_3_text_1': 'Проспективное когортное исследование (2021-2024) с анализом 12,583 амбулаторных визитов 1,070 пациентов. Дополнительно — чувствительный анализ с перекрестным дизайном "case-crossover" для учета временных факторов.',
            'slide_3_text_2': '1,070 пациентов с ревматоидным артритом в крупном медицинском центре Южной Кореи.',
            'slide_3_text_3': 'Оценка 6 загрязнителей воздуха: SO2, NO2, O3, CO, PM10, PM2.5 (ежемесячные уровни) Оценка активности заболевания и обострений при каждом визите Учет демографических, серологических, медикаментозных, социально-экономических и погодных факторов',
            'slide_3_text_4': 'Сбор данных в реальных клинических условиях в течение 4 лет. Сопоставление уровней загрязнения с показателями активности заболевания.',
            'slide_3_text_5': 'Условная логистическая регрессия для сравнения изменений внутри отдельных пациентов с учетом ежедневных концентраций загрязнителей перед каждым визитом.',

            'slide_4_text_1': 'Более высокая концентрация PM2.5 связана с повышенной активностью заболевания и риском обострений. Наиболее выраженный эффект наблюдался при длительном воздействии (более 2 недель) высоких уровней PM2.5.',
            'slide_4_text_1': 'Среди всех изученных загрязнителей (SO2, NO2, O3, CO, PM10, PM2.5) именно PM2.5 был основным фактором, связанным с активностью РА Частицы PM2.5 меньше эритроцитов и могут проникать из легких в кровоток, распространяясь по всему организму',
            'slide_4_text_1': 'Исследование использовало чувствительный анализ для минимизации влияния постоянных факторов и исключения возможности, что активность заболевания влияла на измерения загрязнения. Статистические методы включали условную логистическую регрессию.',
            'slide_4_text_1': 'Не указаны в статье.',

            'slide_5_text_1': 'Воздействие загрязненного воздуха, особенно PM2.5, усиливает активность ревматоидного артрита и увеличивает риск болезненных обострений. Возможный механизм — стимуляция чрезмерной продукции активных форм кислорода, повреждение ДНК и активация воспалительных реакций.',
            'slide_5_text_2': 'Результаты имеют важные клинические, биологические и общественно-здравоохранительные последствия Подтверждают, что вдыхаемые вещества влияют на риск и прогрессирование РА и других аутоиммунных заболеваний Могут помочь врачам объяснить непредсказуемые обострения РА Важны для разработки политики в области общественного здравоохранения',
            'slide_5_text_3': 'Исследование проведено на корейской популяции с определенным генетическим фоном и в специфических экологических условиях Требуются дальнейшие исследования для подтверждения данных в других регионах мира Неизвестно, может ли улучшение качества воздуха напрямую снизить активность заболевания',
            'slide_5_text_4': 'Необходимо определить, может ли снижение воздействия загрязнения или улучшение качества воздуха напрямую уменьшить активность заболевания у пациентов с РА. Также требуется проверка результатов в других популяциях.',
            'slide_5_text_5': 'Длительное воздействие мелкодисперсных частиц PM2.5 в загрязненном воздухе повышает активность ревматоидного артрита и риск обострений, поэтому пациентам с РА рекомендуется избегать продолжительного нахождения в условиях плохого качества воздуха.',
        }

        match _pattern_name:
            case 'pattern_1':
                return PresentationData(
                    name = _data.get('presentation_name'),
                    slides = [

                        # Титульный слайд
                        _Slide(
                            layout=_SlideLayout.BLANK,
                            objects=[

                                _SlideObject(
                                    type=_ObjectType.IMAGE,
                                    content=self._load_svg('../resources/Presentations/Templates/Pattern_1/Backgrounds/title_slide.svg'),
                                    position=(0, 0),
                                    size=(1920, 1080),
                                    text_style=None,
                                    z_index=0
                                ),

                                _SlideObject(
                                    type=_ObjectType.TITLE,
                                    content='Заголовок презентации',
                                    position=(300, 300),
                                    size=(700, 200),
                                    vertically_centered=True,
                                    horizontally_centred=True,
                        
                                    text_style=_ObjectTextStyle(
                                        font_size=50,
                                        font_name='Arial',
                                        bold=False,
                                        italic=False,
                                        color='#FFFFFF',
                                        alignment='right'
                                    ),
                                    z_index=0
                                )
                            ]
                        )
                    ]
                )




































logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s | %(levelname)-8s | %(name)-20s | %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
for logger_name in ['PIL', 'PIL.PngImagePlugin', 'svglib', 'svglib.svglib', 'reportlab']:
    logging.getLogger(logger_name).setLevel(logging.WARNING)


settings = {}
manager = PresentationManager(settings)
manager.create()



